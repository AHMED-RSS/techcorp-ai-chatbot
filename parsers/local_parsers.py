from __future__ import annotations

import csv
import io
import json
import re
import zipfile

from pathlib import Path
from typing import Any
from xml.etree import ElementTree


from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation


from parsers.models import ParsedFile
from parsers.registry import ParserRegistry



TEXT_EXTENSIONS = {
    "txt",
    "md",
    "markdown",
    "rst",
    "log",
    "ini",
    "cfg",
    "conf",
    "toml",
    "yaml",
    "yml",
}



CODE_EXTENSIONS = {
    "py",
    "js",
    "jsx",
    "ts",
    "tsx",
    "java",
    "c",
    "cpp",
    "h",
    "hpp",
    "cs",
    "go",
    "rs",
    "php",
    "rb",
    "swift",
    "kt",
    "kts",
    "scala",
    "sh",
    "bash",
    "ps1",
    "bat",
    "cmd",
    "html",
    "css",
    "scss",
    "sql",
}



IMAGE_EXTENSIONS = {
    "png",
    "jpg",
    "jpeg",
    "webp",
    "bmp",
    "gif",
}



SAFE_ARCHIVE_TEXT_EXTENSIONS = (

    TEXT_EXTENSIONS

    | CODE_EXTENSIONS

    | {

        "csv",

        "json",

        "xml",

    }

)





def clean_title(
    filename: str,
) -> str:


    stem = Path(filename).stem


    cleaned = re.sub(
        r"[_-]+",
        " ",
        stem,
    )


    cleaned = re.sub(
        r"\s+",
        " ",
        cleaned,
    ).strip()


    return cleaned or filename





def decode_text_bytes(
    data: bytes,
) -> tuple[str, str]:


    encodings = (

        "utf-8",

        "utf-8-sig",

        "utf-16",

        "cp1252",

        "latin-1",

    )


    for encoding in encodings:


        try:

            return (

                data.decode(
                    encoding
                ),

                encoding,

            )


        except UnicodeDecodeError:

            continue



    return (

        data.decode(

            "utf-8",

            errors="replace",

        ),

        "utf-8-replacement",

    )





def parse_plain_text(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    extension = (

        Path(filename)

        .suffix

        .lower()

        .lstrip(".")

    )


    return ParsedFile(

        title=clean_title(
            filename
        ),

        text=text,

        category="text",

        mime_type=mime_type,

        extension=extension,

        metadata={

            "encoding":

                encoding,

            "line_count":

                len(
                    text.splitlines()
                ),

        },

    )





def parse_source_code(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    extension = (

        Path(filename)

        .suffix

        .lower()

        .lstrip(".")

    )


    return ParsedFile(

        title=filename,

        text=text,

        category="code",

        mime_type=mime_type,

        extension=extension,

        metadata={

            "language_hint":

                extension,

            "encoding":

                encoding,

            "line_count":

                len(
                    text.splitlines()
                ),

        },

    )
def parse_pdf(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    reader = PdfReader(

        io.BytesIO(data)

    )


    pages: list[str] = []

    warnings: list[str] = []



    for page_number, page in enumerate(

        reader.pages,

        start=1

    ):


        try:

            page_text = (

                page.extract_text()

                or ""

            ).strip()



        except Exception as exc:


            warnings.append(

                f"Page {page_number} extraction failed: {exc}"

            )


            page_text = ""



        if page_text:


            pages.append(

                f"[Page {page_number}]\n"

                f"{page_text}"

            )



    text = "\n\n".join(

        pages

    )



    if not text:


        warnings.append(

            "No text extracted from PDF."

        )



    metadata: dict[str, Any] = {


        "page_count":

            len(reader.pages),


        "ocr_enabled":

            False,

    }



    if reader.metadata:


        metadata.update(

            {

                "pdf_title":

                    reader.metadata.title,


                "author":

                    reader.metadata.author,


                "subject":

                    reader.metadata.subject,

            }

        )



    return ParsedFile(

        title=(

            str(

                metadata.get(

                    "pdf_title"

                )

                or ""

            ).strip()

            or clean_title(

                filename

            )

        ),

        text=text,

        category="document",

        mime_type=mime_type,

        extension="pdf",

        metadata=metadata,

        warnings=warnings,

    )





def parse_docx(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    document = Document(

        io.BytesIO(data)

    )


    parts: list[str] = []



    for paragraph in document.paragraphs:


        value = paragraph.text.strip()


        if value:


            parts.append(value)



    table_count = 0



    for table in document.tables:


        table_count += 1



        for row in table.rows:


            values = [

                cell.text.strip()

                for cell in row.cells

            ]


            if any(values):


                parts.append(

                    " | ".join(values)

                )



    return ParsedFile(

        title=clean_title(filename),

        text="\n\n".join(parts),

        category="document",

        mime_type=mime_type,

        extension="docx",

        metadata={

            "paragraph_count":

                len(document.paragraphs),

            "table_count":

                table_count,

        },

    )





def parse_csv_file(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    rows = []



    try:

        reader = csv.reader(

            io.StringIO(text)

        )


        for row in reader:

            rows.append(

                " | ".join(row)

            )


    except Exception:

        rows = [

            text

        ]



    return ParsedFile(

        title=clean_title(filename),

        text="\n".join(rows),

        category="data",

        mime_type=mime_type,

        extension="csv",

        metadata={

            "encoding":

                encoding,

            "rows":

                len(rows),

        },

    )





def parse_json_file(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    try:

        obj = json.loads(

            text

        )


        text = json.dumps(

            obj,

            indent=2,

            ensure_ascii=False,

        )


    except Exception:

        pass



    return ParsedFile(

        title=clean_title(filename),

        text=text,

        category="data",

        mime_type=mime_type,

        extension="json",

        metadata={

            "encoding":

                encoding,

        },

    )
def parse_xml_file(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    try:

        root = ElementTree.fromstring(
            text
        )


        text = "\n".join(

            root.itertext()

        )


    except Exception:

        pass



    return ParsedFile(

        title=clean_title(filename),

        text=text,

        category="data",

        mime_type=mime_type,

        extension="xml",

        metadata={

            "encoding":

                encoding,

        },

    )





def parse_xlsx(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    workbook = load_workbook(

        io.BytesIO(data),

        read_only=True,

        data_only=True,

    )


    output = []



    for sheet in workbook.worksheets:


        output.append(

            f"Sheet: {sheet.title}"

        )


        for row in sheet.iter_rows(

            values_only=True

        ):


            values = [

                str(value)

                for value in row

                if value is not None

            ]


            if values:


                output.append(

                    " | ".join(values)

                )



    return ParsedFile(

        title=clean_title(filename),

        text="\n".join(output),

        category="data",

        mime_type=mime_type,

        extension="xlsx",

        metadata={

            "sheet_count":

                len(workbook.sheetnames)

        },

    )





def parse_pptx(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    presentation = Presentation(

        io.BytesIO(data)

    )


    slides = []



    for index, slide in enumerate(

        presentation.slides,

        start=1

    ):


        texts = []



        for shape in slide.shapes:


            if hasattr(shape, "text"):


                value = shape.text.strip()



                if value:

                    texts.append(value)



        if texts:


            slides.append(

                f"[Slide {index}]\n"

                + "\n".join(texts)

            )



    return ParsedFile(

        title=clean_title(filename),

        text="\n\n".join(slides),

        category="document",

        mime_type=mime_type,

        extension="pptx",

        metadata={

            "slide_count":

                len(presentation.slides)

        },

    )





def parse_image(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    try:


        image = Image.open(

            io.BytesIO(data)

        )


        metadata = {

            "width":

                image.width,

            "height":

                image.height,

        }



    except Exception as exc:


        metadata = {

            "error":

                str(exc)

        }



    return ParsedFile(

        title=clean_title(filename),

        text="",

        category="image",

        mime_type=mime_type,

        extension=(

            Path(filename)

            .suffix

            .lower()

            .lstrip(".")

        ),

        metadata=metadata,

    )





def parse_zip_archive(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    parts = []

    warnings = []



    try:


        with zipfile.ZipFile(

            io.BytesIO(data)

        ) as archive:



            for name in archive.namelist():


                extension = (

                    Path(name)

                    .suffix

                    .lower()

                    .lstrip(".")

                )



                if extension in SAFE_ARCHIVE_TEXT_EXTENSIONS:


                    try:


                        content = archive.read(

                            name

                        )


                        text, _ = decode_text_bytes(

                            content

                        )


                        parts.append(

                            f"File: {name}\n{text}"

                        )



                    except Exception as exc:


                        warnings.append(

                            f"{name}: {exc}"

                        )



    except Exception as exc:


        warnings.append(

            str(exc)

        )



    return ParsedFile(

        title=clean_title(filename),

        text="\n\n".join(parts),

        category="archive",

        mime_type=mime_type,

        extension="zip",

        metadata={

            "files_found":

                len(parts)

        },

        warnings=warnings,

    )
def parse_unknown_file(
    data: bytes,
    filename: str,
    mime_type: str,
) -> ParsedFile:


    text, encoding = decode_text_bytes(
        data
    )


    return ParsedFile(

        title=clean_title(filename),

        text=text,

        category="unknown",

        mime_type=mime_type,

        extension=(

            Path(filename)

            .suffix

            .lower()

            .lstrip(".")

        ),

        metadata={

            "encoding":

                encoding,

            "fallback_parser":

                True,

        },

    )





def build_default_parser_registry() -> ParserRegistry:

    registry = ParserRegistry()


    registry.register(
        extensions=tuple(TEXT_EXTENSIONS),
        mime_types=(
            "text/plain",
            "text/markdown",
        ),
        parser=parse_plain_text,
        category="text",
    )


    registry.register(
        extensions=tuple(CODE_EXTENSIONS),
        mime_types=(
            "text/plain",
            "application/octet-stream",
        ),
        parser=parse_source_code,
        category="code",
    )


    registry.register(
        extensions=("pdf",),
        mime_types=("application/pdf",),
        parser=parse_pdf,
        category="document",
    )


    registry.register(
        extensions=("docx",),
        mime_types=(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ),
        parser=parse_docx,
        category="document",
    )


    registry.register(
        extensions=("csv",),
        mime_types=("text/csv",),
        parser=parse_csv_file,
        category="data",
    )


    registry.register(
        extensions=("json",),
        mime_types=("application/json",),
        parser=parse_json_file,
        category="data",
    )


    registry.register(
        extensions=("xml",),
        mime_types=(
            "application/xml",
            "text/xml",
        ),
        parser=parse_xml_file,
        category="data",
    )


    registry.register(
        extensions=("xlsx",),
        mime_types=(
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ),
        parser=parse_xlsx,
        category="data",
    )


    registry.register(
        extensions=("pptx",),
        mime_types=(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ),
        parser=parse_pptx,
        category="document",
    )


    registry.register(
        extensions=tuple(IMAGE_EXTENSIONS),
        mime_types=("image/*",),
        parser=parse_image,
        category="image",
    )


    registry.register(
        extensions=("zip",),
        mime_types=("application/zip",),
        parser=parse_zip_archive,
        category="archive",
    )


    return registry